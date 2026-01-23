from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # 定義輔助函式：新增標題與內容頁
    def add_slide(title_text, content_text_list):
        slide_layout = prs.slide_layouts[1] # 1 是標題+內容版型
        slide = prs.slides.add_slide(slide_layout)
        
        # 設定標題
        title = slide.shapes.title
        title.text = title_text
        
        # 設定內容
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        
        for i, text in enumerate(content_text_list):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(24)
            # 簡單的層級處理 (如果文字開頭是 -, 縮排)
            if text.startswith("  -"):
                p.level = 1
                p.text = text.replace("  -", "").strip()

    # --- 第 1 頁：封面 ---
    slide_layout = prs.slide_layouts[0] # 0 是封面版型
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    title.text = "下一代自動化測試框架\nExcel 驅動與 AI 視覺整合"
    subtitle.text = "基於 Python 的多平台 (Desktop/Web/Mobile) 解決方案\n\n報告人：[您的名字]"

    # --- 第 2 頁：痛點與目標 ---
    add_slide("面臨挑戰與設計目標", [
        "痛點 1：測試維護成本高",
        "  - 傳統腳本寫死在代碼中，需求變更需大量修改",
        "痛點 2：桌面端自動化不穩定",
        "  - 傳統座標/圖像比對對解析度與渲染差異容忍度低",
        "痛點 3：技術門檻高",
        "  - 非技術人員 (QA/PM) 難以參與維護",
        "",
        "設計目標：",
        "  - 數據驅動 (Data-Driven)",
        "  - 高魯棒性 (Robustness, AI 輔助)",
        "  - 低門檻 (Excel + GUI)"
    ])

    # --- 第 3 頁：核心架構 ---
    add_slide("核心架構 (Architecture)", [
        "1. 測試定義層 (Definition Layer)",
        "  - Excel (TestPlan.xlsx): 定義步驟與參數",
        "2. 執行引擎層 (Engine Layer)",
        "  - TestRunner: 調度核心 (Pytest)",
        "  - StepTranslator: 關鍵字翻譯器",
        "3. 業務邏輯層 (Action Layer)",
        "  - NxPocActions, NxMobileActions",
        "4. 基礎操作層 (Base Layer)",
        "  - Smart Click Engine (AI/OCR/CV)"
    ])

    # --- 第 4 頁：技術亮點 I ---
    add_slide("技術亮點 I - Excel 關鍵字驅動", [
        "實現「寫 Excel 即寫腳本」",
        "",
        "運作原理：",
        "  - 讀取 FlowName 與 Params",
        "  - 透過反射機制 (Reflection) 動態呼叫 Python",
        "優勢：",
        "  - 新增測試案例無需修改程式碼",
        "  - 測試資產可由非開發人員維護"
    ])

    # --- 第 5 頁：技術亮點 II ---
    add_slide("技術亮點 II - 智能混合辨識 (Smart Click)", [
        "解決「定位失敗」的多層級保底策略：",
        "",
        "Level 1: AI 視覺理解 (VLM)",
        "  - 使用本地大模型 (Ollama/LLaVA) 理解自然語言",
        "Level 2: 圖像特徵比對 (OpenCV)",
        "Level 3: OCR 文字辨識 (PaddleOCR)",
        "Level 4: 自學習座標保底 (Coordinate Fallback)"
    ])

    # --- 第 6 頁：技術亮點 III ---
    add_slide("技術亮點 III - 可視化報告與除錯", [
        "HTML 互動式報告 (類似 Allure/UFT)：",
        "",
        "智慧標註截圖：",
        "  - 自動繪製紅框/綠框，標示 AI 識別位置",
        "  - 展示「預期點擊」與「實際識別」差異",
        "失敗現場還原：",
        "  - 自動保存 Terminal Log 與全螢幕截圖"
    ])

    # --- 第 7 頁：開發者體驗 ---
    add_slide("開發者體驗 - GUI 啟動器", [
        "讓自動化觸手可及",
        "",
        "Tkinter GUI 工具 (test_case_launcher.py)：",
        "  - 自動掃描 Excel 測試計畫",
        "  - 勾選介面，支援批量執行",
        "  - 多執行緒設計，介面不卡頓",
        "部署優勢：",
        "  - 可打包為 .exe，無 Python 環境也能執行"
    ])

    # --- 第 8 頁：總結 ---
    add_slide("總結與成效", [
        "跨平台整合：",
        "  - Desktop, Web, Mobile 統一架構",
        "穩定性提升：",
        "  - AI 輔助大幅降低 Flaky Tests",
        "維護效率：",
        "  - 修改時間從「小時級」縮短至「分鐘級」",
        "",
        "未來展望：接入 GPT-4o 提升辨識速度"
    ])

    # 存檔
    output_file = "Automation_Framework_Presentation.pptx"
    prs.save(output_file)
    print(f"✅ 簡報已生成！請開啟：{output_file}")

if __name__ == "__main__":
    create_presentation()